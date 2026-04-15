# data_lake_ingestor.py
# Script ETL completo para convertir archivos heterogéneos en un data lake analítico en formato Parquet.
#
# Dependencias requeridas (ver requirements_ingestor.txt):
#   pyarrow>=14.0
#   polars>=0.20
#   puremagic>=1.28
#   pdfplumber>=0.10
#   python-docx>=1.1
#   python-pptx>=0.6
#   beautifulsoup4>=4.12
#   lxml>=5.0
#   pillow>=10.0
#   tqdm>=4.66
#   openpyxl>=3.1

import hashlib
import json
import logging
import os
import io
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional

import puremagic as magic
import polars as pl
import pyarrow as pa
import pyarrow.ipc as ipc
import pyarrow.parquet as pq
from tqdm import tqdm

# ================= CONFIGURACIÓN =================
DIRECTORIO_RAIZ = Path("F:/Sistema de gestion seguridad")
ARCHIVO_LOG = Path("F:/procesados_log.json")
ARCHIVO_DATASET = Path("F:/dataset_maestro_seguridad.parquet")
MAX_WORKERS = 4
LIMITE_MEMORIA_ARCHIVO = 200 * 1024 * 1024  # 200 MB

# ================= LOGGING =================
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
)

# ================= CONSTANTES =================
NOMBRES_IGNORADOS = {"desktop.ini", "thumbs.db", ".ds_store"}

# Schema explícito de PyArrow para el dataset final
SCHEMA_ARROW = pa.schema([
    pa.field("ruta", pa.string()),
    pa.field("nombre", pa.string()),
    pa.field("hash", pa.string()),
    pa.field("mtime", pa.timestamp("us")),
    pa.field("tamano_bytes", pa.int64()),
    pa.field("tipo_mime", pa.string()),
    pa.field("extension", pa.string()),
    pa.field("contenido_texto", pa.string()),
    pa.field("metadata_json", pa.string()),
    pa.field("tabla_arrow", pa.binary()),
    pa.field("fecha_procesado", pa.timestamp("us")),
])


# ================= UTILIDADES =================

def calcular_hash_seguro(ruta: Path) -> str:
    """Calcula SHA-256 del archivo por bloques de 4 KB para no saturar la RAM."""
    sha256 = hashlib.sha256()
    with open(ruta, "rb") as f:
        for bloque in iter(lambda: f.read(4096), b""):
            sha256.update(bloque)
    return sha256.hexdigest()


def tabla_a_bytes_arrow(tabla: pa.Table) -> bytes:
    """Serializa una tabla Arrow al formato IPC en memoria."""
    sink = pa.BufferOutputStream()
    with ipc.new_stream(sink, tabla.schema) as writer:
        writer.write_table(tabla)
    return sink.getvalue().to_pybytes()


def detectar_mime(ruta: Path) -> str:
    """Detecta el tipo MIME real del archivo usando puremagic (sin DLLs externas)."""
    try:
        coincidencias = magic.from_file(str(ruta))
        return coincidencias[0].mime_type if coincidencias else "application/octet-stream"
    except Exception as e:
        logging.error(f"Error detectando MIME en {ruta}: {e}")
        return "application/octet-stream"


# ================= HANDLERS =================
handlers: Dict[str, Any] = {}


def registrar_handler(extensiones: List[str]):
    """Decorador para registrar handlers por extensión de archivo."""
    def decorador(func):
        for ext in extensiones:
            handlers[ext.lower()] = func
        return func
    return decorador


@registrar_handler([".txt", ".log", ".md", ".py", ".js", ".json"])
def manejar_texto(ruta: Path, mime: str) -> Dict[str, Any]:
    """Lee archivos de texto plano en UTF-8."""
    try:
        texto = ruta.read_text(encoding="utf-8", errors="ignore")
        return {
            "contenido_texto": texto,
            "tabla_arrow": None,
            "metadata": {"lineas": len(texto.splitlines())},
        }
    except Exception as e:
        logging.error(f"Error leyendo texto en {ruta}: {e}")
        return {"contenido_texto": "", "tabla_arrow": None, "metadata": {}}


@registrar_handler([".csv"])
def manejar_csv(ruta: Path, mime: str) -> Dict[str, Any]:
    """Lee CSV con Polars (modo lazy) y serializa la tabla a Arrow IPC."""
    try:
        df = pl.scan_csv(str(ruta), ignore_errors=True).collect()
        tabla = df.to_arrow()
        return {
            "contenido_texto": f"CSV con {tabla.num_rows} filas y {len(tabla.column_names)} columnas",
            "tabla_arrow": tabla_a_bytes_arrow(tabla),
            "metadata": {"filas": tabla.num_rows, "columnas": tabla.column_names},
        }
    except Exception as e:
        logging.error(f"Error procesando CSV {ruta}: {e}")
        return {"contenido_texto": "", "tabla_arrow": None, "metadata": {}}


@registrar_handler([".xlsx", ".xls"])
def manejar_excel(ruta: Path, mime: str) -> Dict[str, Any]:
    """Lee Excel con Polars + openpyxl y serializa la tabla a Arrow IPC."""
    try:
        df = pl.read_excel(str(ruta), engine="openpyxl")
        tabla = df.to_arrow()
        return {
            "contenido_texto": f"Excel con {tabla.num_rows} filas y {len(tabla.column_names)} columnas",
            "tabla_arrow": tabla_a_bytes_arrow(tabla),
            "metadata": {"filas": tabla.num_rows, "columnas": tabla.column_names},
        }
    except Exception as e:
        logging.error(f"Error procesando Excel {ruta}: {e}")
        return {"contenido_texto": "", "tabla_arrow": None, "metadata": {}}


@registrar_handler([".pdf"])
def manejar_pdf(ruta: Path, mime: str) -> Dict[str, Any]:
    """Extrae texto de PDF página a página con pdfplumber."""
    try:
        import pdfplumber
        with pdfplumber.open(ruta) as pdf:
            paginas = len(pdf.pages)
            texto = "\n".join(p.extract_text() or "" for p in pdf.pages)
        return {
            "contenido_texto": texto,
            "tabla_arrow": None,
            "metadata": {"paginas": paginas},
        }
    except Exception as e:
        logging.error(f"Error procesando PDF {ruta}: {e}")
        return {"contenido_texto": "", "tabla_arrow": None, "metadata": {}}


@registrar_handler([".docx"])
def manejar_docx(ruta: Path, mime: str) -> Dict[str, Any]:
    """Extrae párrafos de documentos Word con python-docx."""
    try:
        from docx import Document
        doc = Document(ruta)
        texto = "\n".join(p.text for p in doc.paragraphs)
        return {
            "contenido_texto": texto,
            "tabla_arrow": None,
            "metadata": {"parrafos": len(doc.paragraphs)},
        }
    except Exception as e:
        logging.error(f"Error procesando DOCX {ruta}: {e}")
        return {"contenido_texto": "", "tabla_arrow": None, "metadata": {}}


@registrar_handler([".pptx"])
def manejar_pptx(ruta: Path, mime: str) -> Dict[str, Any]:
    """Extrae texto de presentaciones PowerPoint con python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(ruta)
        partes = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    partes.append(shape.text)
        texto = "\n".join(partes)
        return {
            "contenido_texto": texto,
            "tabla_arrow": None,
            "metadata": {"diapositivas": len(prs.slides)},
        }
    except Exception as e:
        logging.error(f"Error procesando PPTX {ruta}: {e}")
        return {"contenido_texto": "", "tabla_arrow": None, "metadata": {}}


@registrar_handler([".html", ".htm"])
def manejar_html(ruta: Path, mime: str) -> Dict[str, Any]:
    """Extrae texto limpio de HTML con BeautifulSoup."""
    try:
        from bs4 import BeautifulSoup
        html = ruta.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "lxml")
        texto = soup.get_text(separator="\n", strip=True)
        return {
            "contenido_texto": texto,
            "tabla_arrow": None,
            "metadata": {"etiquetas": len(soup.find_all())},
        }
    except Exception as e:
        logging.error(f"Error procesando HTML {ruta}: {e}")
        return {"contenido_texto": "", "tabla_arrow": None, "metadata": {}}


@registrar_handler([".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff"])
def manejar_imagen(ruta: Path, mime: str) -> Dict[str, Any]:
    """Extrae metadata EXIF de imágenes con Pillow."""
    try:
        from PIL import Image
        with Image.open(ruta) as img:
            meta = {
                "tamaño": img.size,
                "modo": img.mode,
                "formato": img.format,
            }
            exif_data = img.getexif() if hasattr(img, "getexif") else None
            if exif_data:
                meta["exif_tags"] = len(exif_data)
            contenido = f"Imagen {img.format}: {img.size[0]}x{img.size[1]} px, modo {img.mode}"
        return {
            "contenido_texto": contenido,
            "tabla_arrow": None,
            "metadata": meta,
        }
    except Exception as e:
        logging.error(f"Error procesando imagen {ruta}: {e}")
        return {"contenido_texto": "", "tabla_arrow": None, "metadata": {}}


def manejar_binario(ruta: Path, mime: str) -> Dict[str, Any]:
    """Handler de fallback para extensiones no reconocidas."""
    return {
        "contenido_texto": f"Archivo binario registrado ({mime})",
        "tabla_arrow": None,
        "metadata": {"mime": mime},
    }


# ================= NÚCLEO DEL PROCESAMIENTO =================

def procesar_archivo(ruta: Path) -> Optional[Dict[str, Any]]:
    """
    Procesa un archivo individual y devuelve un diccionario con sus metadatos y contenido.
    Retorna None si el archivo debe omitirse o si ocurre un error irrecuperable.
    """
    try:
        stat = ruta.stat()
        tamano = stat.st_size
        mtime_ts = stat.st_mtime

        # Filtrar archivos temporales, de sistema y vacíos
        if ruta.name.startswith("~$") or ruta.name.lower() in NOMBRES_IGNORADOS or tamano == 0:
            return None

        h = calcular_hash_seguro(ruta)
        ext = ruta.suffix.lower()
        m = detectar_mime(ruta)

        handler = handlers.get(ext, manejar_binario)
        try:
            resultado = handler(ruta, m)
        except Exception as e:
            logging.error(f"Error en handler '{ext}' para {ruta}: {e}")
            resultado = {"contenido_texto": "", "tabla_arrow": None, "metadata": {}}

        return {
            "ruta": str(ruta.absolute()),
            "nombre": ruta.name,
            "hash": h,
            "mtime": datetime.fromtimestamp(mtime_ts, tz=timezone.utc).replace(tzinfo=None),
            "tamano_bytes": tamano,
            "tipo_mime": m,
            "extension": ext,
            "contenido_texto": resultado.get("contenido_texto", "") or "",
            "metadata_json": json.dumps(resultado.get("metadata", {}), ensure_ascii=False, default=str),
            "tabla_arrow": resultado.get("tabla_arrow", None),
            "fecha_procesado": datetime.now(tz=timezone.utc).replace(tzinfo=None),
        }
    except Exception as e:
        logging.error(f"Error procesando {ruta}: {e}")
        return None


def construir_dataframe(resultados: List[Dict[str, Any]]) -> pl.DataFrame:
    """
    Construye un DataFrame Polars con schema explícito a partir de los resultados procesados.
    La columna tabla_arrow se convierte explícitamente para manejar mezcla de bytes y None.
    """
    # Separar columnas para control de tipos
    tabla_arrow_col = [r["tabla_arrow"] for r in resultados]

    # Crear tabla Arrow con schema explícito para evitar inferencia incorrecta de bytes|None
    arrays = {
        "ruta": pa.array([r["ruta"] for r in resultados], type=pa.string()),
        "nombre": pa.array([r["nombre"] for r in resultados], type=pa.string()),
        "hash": pa.array([r["hash"] for r in resultados], type=pa.string()),
        "mtime": pa.array([r["mtime"] for r in resultados], type=pa.timestamp("us")),
        "tamano_bytes": pa.array([r["tamano_bytes"] for r in resultados], type=pa.int64()),
        "tipo_mime": pa.array([r["tipo_mime"] for r in resultados], type=pa.string()),
        "extension": pa.array([r["extension"] for r in resultados], type=pa.string()),
        "contenido_texto": pa.array([r["contenido_texto"] for r in resultados], type=pa.string()),
        "metadata_json": pa.array([r["metadata_json"] for r in resultados], type=pa.string()),
        "tabla_arrow": pa.array(tabla_arrow_col, type=pa.binary()),
        "fecha_procesado": pa.array([r["fecha_procesado"] for r in resultados], type=pa.timestamp("us")),
    }

    tabla_pa = pa.table(arrays, schema=SCHEMA_ARROW)
    return pl.from_arrow(tabla_pa)


# ================= FLUJO PRINCIPAL =================

def ejecutar_mapeo() -> None:
    """
    Flujo principal ETL:
    1. Verifica el directorio raíz
    2. Escanea archivos recursivamente
    3. Filtra ya procesados por mtime
    4. Procesa en paralelo (pequeños) y secuencial (grandes)
    5. Guarda dataset Parquet y log JSON
    """
    if not DIRECTORIO_RAIZ.exists():
        logging.error(f"La ruta {DIRECTORIO_RAIZ} no existe.")
        print(f"❌ La ruta {DIRECTORIO_RAIZ} no existe.")
        return

    print("🔍 Escaneando archivos...")
    todos = [p for p in DIRECTORIO_RAIZ.rglob("*") if p.is_file()]
    logging.info(f"Archivos encontrados: {len(todos)}")

    # Cargar log de archivos ya procesados
    log_actual: Dict[str, str] = {}
    if ARCHIVO_LOG.exists():
        try:
            with open(ARCHIVO_LOG, "r", encoding="utf-8") as f:
                log_actual = json.load(f)
        except Exception as e:
            logging.error(f"Error cargando log {ARCHIVO_LOG}: {e}")

    # Filtrar pendientes por mtime
    pendientes: List[Path] = []
    for p in todos:
        try:
            mtime_actual = str(p.stat().st_mtime)
            if log_actual.get(str(p)) == mtime_actual:
                continue
            pendientes.append(p)
        except Exception as e:
            logging.error(f"Error accediendo a stats de {p}: {e}")

    if not pendientes:
        print("✅ Todo actualizado. No hay archivos nuevos o modificados.")
        return

    print(f"📂 {len(pendientes)} archivo(s) nuevo(s) o modificado(s) a procesar.")

    # Separar en pequeños (<= 200MB) y grandes (> 200MB)
    pequenos = [p for p in pendientes if p.stat().st_size <= LIMITE_MEMORIA_ARCHIVO]
    grandes = [p for p in pendientes if p.stat().st_size > LIMITE_MEMORIA_ARCHIVO]

    logging.info(f"Archivos pequeños (≤200MB): {len(pequenos)} | Archivos grandes (>200MB): {len(grandes)}")

    resultados: List[Dict[str, Any]] = []
    errores = 0

    # Procesar archivos pequeños en paralelo
    if pequenos:
        print(f"⚙️  Procesando {len(pequenos)} archivo(s) pequeño(s) en paralelo (workers={MAX_WORKERS})...")
        with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
            futuros = {executor.submit(procesar_archivo, p): p for p in pequenos}
            for futuro in tqdm(as_completed(futuros), total=len(pequenos), desc="Pequeños"):
                try:
                    res = futuro.result()
                    if res:
                        resultados.append(res)
                        log_actual[res["ruta"]] = str(os.path.getmtime(res["ruta"]))
                    else:
                        errores += 1
                except Exception as e:
                    ruta_err = futuros[futuro]
                    logging.error(f"Error obteniendo resultado de {ruta_err}: {e}")
                    errores += 1

    # Procesar archivos grandes de forma secuencial
    if grandes:
        print(f"🐘 Procesando {len(grandes)} archivo(s) grande(s) de forma secuencial...")
        for p in tqdm(grandes, desc="Grandes"):
            res = procesar_archivo(p)
            if res:
                resultados.append(res)
                log_actual[res["ruta"]] = str(os.path.getmtime(res["ruta"]))
            else:
                errores += 1

    if not resultados:
        print("⚠️  No se generaron registros válidos en este ciclo.")
        logging.warning("No se generaron registros válidos.")
        return

    # Construir DataFrame nuevo con schema explícito
    df_nuevo = construir_dataframe(resultados)

    # Combinar con dataset existente y deduplicar
    if ARCHIVO_DATASET.exists():
        try:
            df_viejo = pl.read_parquet(ARCHIVO_DATASET)
            df_final = (
                pl.concat([df_viejo, df_nuevo], how="diagonal")
                .sort("fecha_procesado")
                .unique(subset=["hash"], keep="last")
            )
        except Exception as e:
            logging.error(f"Error cargando dataset existente {ARCHIVO_DATASET}: {e}")
            df_final = df_nuevo
    else:
        df_final = df_nuevo

    # Guardar dataset Parquet
    try:
        ARCHIVO_DATASET.parent.mkdir(parents=True, exist_ok=True)
        df_final.write_parquet(ARCHIVO_DATASET)
        logging.info(f"Dataset guardado en {ARCHIVO_DATASET} ({len(df_final)} registros totales)")
    except Exception as e:
        logging.error(f"Error guardando dataset Parquet {ARCHIVO_DATASET}: {e}")

    # Guardar log JSON actualizado
    try:
        ARCHIVO_LOG.parent.mkdir(parents=True, exist_ok=True)
        with open(ARCHIVO_LOG, "w", encoding="utf-8") as f:
            json.dump(log_actual, f, indent=2, ensure_ascii=False)
        logging.info(f"Log guardado en {ARCHIVO_LOG}")
    except Exception as e:
        logging.error(f"Error guardando log {ARCHIVO_LOG}: {e}")

    nuevos = len(resultados)
    totales = len(df_final)
    print(f"\n🎉 Proceso finalizado.")
    print(f"   Registros totales en dataset : {totales}")
    print(f"   Nuevos registros en este ciclo: {nuevos}")
    print(f"   Archivos con error/omitidos   : {errores}")


# ================= PUNTO DE ENTRADA =================
if __name__ == "__main__":
    try:
        ejecutar_mapeo()
    except KeyboardInterrupt:
        print("\n⚠️ Proceso interrumpido por el usuario.")
